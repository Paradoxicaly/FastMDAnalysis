#!/usr/bin/env python3
"""Insert or update scaling benchmark slides in benchmark_presentation.pptx."""
from __future__ import annotations

import argparse
from pathlib import Path
from typing import Iterable, List, Sequence, Tuple

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches

SLIDE_LAYOUT_INDEX = 5  # Title Only
DEFAULT_PPT = Path("benchmark_presentation.pptx")
DEFAULT_IMAGE_DIR = Path("assets/benchmarks")

IMAGE_SPEC: Sequence[Tuple[str, str]] = (
    ("FastMDAnalysis Runtime Scaling", "fastmdanalysis_runtime_scaling.png"),
    ("FastMDAnalysis Memory Scaling", "fastmdanalysis_memory_scaling.png"),
    ("MDTraj Runtime Scaling", "mdtraj_runtime_scaling.png"),
    ("MDTraj Memory Scaling", "mdtraj_memory_scaling.png"),
    ("MDAnalysis Runtime Scaling", "mdanalysis_runtime_scaling.png"),
    ("MDAnalysis Memory Scaling", "mdanalysis_memory_scaling.png"),
    ("Combined Runtime Scaling (Linear)", "combined_runtime_scaling_linear.png"),
    ("Combined Runtime Scaling (Log)", "combined_runtime_scaling_log.png"),
    ("Combined Memory Scaling (Linear)", "combined_memory_scaling_linear.png"),
    ("Combined Memory Scaling (Log)", "combined_memory_scaling_log.png"),
    ("Lines of Code Comparison", "loc_comparison.png"),
)


def _clear_non_title_shapes(slide) -> None:
    title = slide.shapes.title
    for shape in list(slide.shapes):
        if title is not None and shape == title:
            continue
        slide.shapes._spTree.remove(shape._element)


def _ensure_slide(prs: Presentation, title_text: str):
    for slide in prs.slides:
        title = slide.shapes.title
        if title and title.text.strip() == title_text:
            return slide
    layout = prs.slide_layouts[SLIDE_LAYOUT_INDEX]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title_text
    return slide


def insert_images(ppt_path: Path, image_dir: Path, specs: Iterable[Tuple[str, str]]) -> List[Tuple[str, Path]]:
    prs = Presentation(ppt_path)
    added: List[Tuple[str, Path]] = []
    margin = Inches(0.4)
    top_offset = Inches(1.2)
    img_width = prs.slide_width - 2 * margin

    for title_text, filename in specs:
        image_path = image_dir / filename
        if not image_path.exists():
            raise FileNotFoundError(f"Missing image: {image_path}")
        slide = _ensure_slide(prs, title_text)
        _clear_non_title_shapes(slide)
        slide.shapes.add_picture(str(image_path), margin, top_offset, width=img_width)
        added.append((title_text, image_path))

    prs.save(ppt_path)
    return added


def parse_args(argv: Sequence[str] | None = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Embed scaling PNGs into the benchmark presentation.")
    parser.add_argument("--ppt", type=Path, default=DEFAULT_PPT, help="Path to PPTX file to modify.")
    parser.add_argument("--images", type=Path, default=DEFAULT_IMAGE_DIR, help="Directory containing scaling PNGs.")
    return parser.parse_args(argv)


def main(argv: Sequence[str] | None = None) -> None:
    args = parse_args(argv)
    inserted = insert_images(args.ppt, args.images, IMAGE_SPEC)
    print(f"Updated {args.ppt} with {len(inserted)} scaling slide(s):")
    for title, path in inserted:
        print(f"  • {title} <- {path}")


if __name__ == "__main__":
    main()
