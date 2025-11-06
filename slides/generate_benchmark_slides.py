from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Iterable, List, Sequence

from pptx import Presentation
from pptx.util import Inches, Pt
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

AXIS_LABEL_SIZE = 18
TICK_LABEL_SIZE = 16
TITLE_FONT_SIZE = 20
LEGEND_FONT_SIZE = 13
plt.rcParams.update({
    "axes.labelsize": AXIS_LABEL_SIZE,
    "axes.titlesize": TITLE_FONT_SIZE,
    "figure.titlesize": TITLE_FONT_SIZE,
    "legend.fontsize": LEGEND_FONT_SIZE,
    "xtick.labelsize": TICK_LABEL_SIZE,
    "ytick.labelsize": TICK_LABEL_SIZE,
})

ROOT = Path(__file__).resolve().parents[1]
BENCHMARK_DIR = ROOT / "benchmarks"
if str(BENCHMARK_DIR) not in sys.path:
    sys.path.insert(0, str(BENCHMARK_DIR))

from dataset_config import get_dataset_config, list_datasets  # type: ignore[import-error]
RESULTS_ROOT = ROOT / "benchmarks" / "results"

DEFAULT_DATASET = "trpcage"
DATASET_SLUG = DEFAULT_DATASET
DATASET_LABEL = ""
OUTPUT_PPTX = ROOT / "slides" / "FastMDAnalysis_benchmarks.pptx"
OVERVIEW_DIR = RESULTS_ROOT / "overview"
INSTRUMENTATION_JSON = OVERVIEW_DIR / "instrumentation_summary.json"
INSTRUMENTATION_CHART = OVERVIEW_DIR / "instrumentation_overview.png"


def _set_dataset(slug: str, *, output_path: Path | None = None) -> None:
    global DATASET_SLUG, DATASET_LABEL, OUTPUT_PPTX, OVERVIEW_DIR, INSTRUMENTATION_JSON, INSTRUMENTATION_CHART
    config = get_dataset_config(slug)
    DATASET_SLUG = config.slug
    DATASET_LABEL = config.label
    OVERVIEW_DIR = RESULTS_ROOT / f"overview_{DATASET_SLUG}"
    INSTRUMENTATION_JSON = OVERVIEW_DIR / "instrumentation_summary.json"
    INSTRUMENTATION_CHART = OVERVIEW_DIR / "instrumentation_overview.png"
    OUTPUT_PPTX = (
        output_path
        if output_path is not None
        else ROOT / "slides" / f"FastMDAnalysis_benchmarks_{DATASET_SLUG}.pptx"
    )


_set_dataset(DEFAULT_DATASET)

PLOT_FILES = [
    ("Runtime Summary", "runtime_summary.png"),
    ("Runtime Breakdown", "runtime_summary_split.png"),
    ("Runtime Breakdown (All Tools)", "runtime_summary_split_all.png"),
    ("Peak Memory", "memory_summary.png"),
    ("Peak Memory Distribution", "memory_distribution.png"),
    ("Runtime Distribution", "runtime_distribution.png"),
    ("Line-of-Code Summary", "loc_summary.png"),
]

def _overview_charts() -> List[dict[str, object]]:
    return [
        {
            "title": "Modules touched per benchmark",
            "path": OVERVIEW_DIR / "modules_per_benchmark.png",
            "caption": "Distinct modules per workflow; FastMDAnalysis stays flat",
        },
        {
            "title": "Functions touched per benchmark",
            "path": OVERVIEW_DIR / "functions_per_benchmark.png",
            "caption": "Function touchpoints per workload",
        },
        {
            "title": "External module dependencies",
            "path": OVERVIEW_DIR / "external_modules.png",
            "caption": "Non-core imports each tool requires",
            "detail_field": "external_modules",
        },
        {
            "title": "Snippet lines of code (calc vs plot)",
            "path": OVERVIEW_DIR / "loc_totals.png",
            "caption": "Human effort split between calculation and plotting code",
            "summary_key": "loc_totals",
        },
        {
            "title": "Runtime footprint overview",
            "path": OVERVIEW_DIR / "runtime_totals.png",
            "caption": "Total computation vs plotting time across benchmarks",
            "summary_key": "runtime_totals",
        },
        {
            "title": "Peak memory footprint overview",
            "path": OVERVIEW_DIR / "peak_mem_totals.png",
            "caption": "Peak memory split between computation and plotting helpers",
            "summary_key": "peak_mem_totals",
        },
    ]


def iter_benchmark_dirs(root: Path, dataset_slug: str) -> Iterable[Path]:
    suffix = f"_{dataset_slug}"
    for entry in sorted(root.iterdir()):
        if not entry.is_dir():
            continue
        if entry.name.startswith("overview"):
            continue
        # Skip any orchestrator outputs (these can be empty and shouldn't create slides)
        if entry.name.lower().startswith("orchestrator"):
            continue
        if entry.name.endswith(suffix):
            yield entry


def add_benchmark_slide(prs: Presentation, benchmark_dir: Path) -> None:
    title_text = benchmark_dir.name.replace("_", " ").title()
    readme = benchmark_dir / "README.md"
    summary_lines = readme.read_text(encoding="utf-8").splitlines() if readme.exists() else []

    charts = [
        (label, benchmark_dir / filename)
        for label, filename in PLOT_FILES
        if (benchmark_dir / filename).exists()
    ]
    if not charts:
        charts = []

    positions_first = [
        (Inches(0.5), Inches(2.0)),
        (Inches(5.0), Inches(2.0)),
        (Inches(0.5), Inches(4.9)),
        (Inches(5.0), Inches(4.9)),
    ]
    positions_continued = [
        (Inches(0.5), Inches(1.2)),
        (Inches(5.0), Inches(1.2)),
        (Inches(0.5), Inches(4.1)),
        (Inches(5.0), Inches(4.1)),
    ]
    image_width = Inches(4.0)

    for slide_index in range(max(1, (len(charts) + 3) // 4)):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        positions = positions_first if slide_index == 0 else positions_continued

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9.0), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_run = title_frame.paragraphs[0].add_run()
        title_suffix = "" if slide_index == 0 else " (continued)"
        title_run.text = f"{title_text}{title_suffix}"
        title_run.font.bold = True
        title_run.font.size = Pt(28)

        if slide_index == 0:
            summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9.0), Inches(1.6))
            summary_frame = summary_box.text_frame
            summary_frame.clear()
            if summary_lines:
                for idx, line in enumerate(summary_lines[:3]):
                    paragraph = summary_frame.add_paragraph() if idx else summary_frame.paragraphs[0]
                    paragraph.text = line.lstrip("- ")
                    paragraph.font.size = Pt(14)
            note_para = summary_frame.add_paragraph() if summary_lines else summary_frame.paragraphs[0]
            note_para.text = (
                "Note: MDTraj and MDAnalysis metrics include helper scripts we wrote to save data files and plots; "
                "their APIs do not emit those artifacts automatically."
            )
            note_para.font.size = Pt(12)
            note_para.font.italic = True

        start = slide_index * len(positions)
        group = charts[start:start + len(positions)]
        for (_label, image_path), (left, top) in zip(group, positions):
            slide.shapes.add_picture(str(image_path), left, top, width=image_width)


def add_chart_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
    caption: str,
    notes: Sequence[str] | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Titles and captions are deliberately omitted so only the plots are visible on the slide.
    if image_path.exists():
        slide.shapes.add_picture(
            str(image_path),
            Inches(0.5),
            Inches(1.0),
            width=Inches(9.0),
        )

    if notes:
        caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(9.0), Inches(1.6))
        caption_frame = caption_box.text_frame
        caption_frame.clear()
        for idx, note in enumerate(notes):
            paragraph = caption_frame.add_paragraph() if idx else caption_frame.paragraphs[0]
            paragraph.text = note
            paragraph.level = 0
            paragraph.font.size = Pt(12)
            paragraph.font.bold = False


def add_summary_slide(prs: Presentation, summary: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(6.5))
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()

    title_run = frame.paragraphs[0].add_run()
    title_run.text = "Instrumentation Overview"
    title_run.font.size = Pt(28)
    title_run.font.bold = True

    aggregated = summary.get("aggregated", {})
    loc_totals = summary.get("loc_totals", {})
    runtime_totals = summary.get("runtime_totals", {})
    peak_mem_totals = summary.get("peak_mem_totals", {})

    def _fmt_tool(name: str) -> str:
        return "FastMDAnalysis" if name == "fastmdanalysis" else name.capitalize()

    for tool, stats in aggregated.items():
        paragraph = frame.add_paragraph()
        paragraph.level = 1
        paragraph.text = (
            f"{_fmt_tool(tool)}: {stats['modules']} modules, {stats['functions']} functions, "
            f"{stats['external_modules']} external imports, "
            f"{stats['attempts']} runs, {stats['successes']} successful, {stats['exceptions']} exceptions."
        )
        paragraph.font.size = Pt(16)

    if loc_totals:
        snippet_para = frame.add_paragraph()
        snippet_para.level = 1
        snippet_para.font.size = Pt(16)
        chunks = []
        for tool, values in loc_totals.items():
            calc_lines = values.get("calc", 0)
            plot_lines = values.get("plot", 0)
            total_lines = calc_lines + plot_lines
            chunks.append(
                f"{_fmt_tool(tool)} {calc_lines} calc / {plot_lines} plot ({total_lines} total) lines"
            )
        snippet_para.text = "Snippet footprint overview: " + ", ".join(chunks) + "."

    if runtime_totals:
        runtime_para = frame.add_paragraph()
        runtime_para.level = 1
        runtime_para.font.size = Pt(16)
        runtime_chunks = []
        for tool, values in runtime_totals.items():
            calc_runtime = float(values.get("calc", 0.0))
            plot_runtime = float(values.get("plot", 0.0))
            overhead_runtime = float(values.get("overhead", 0.0))
            total_runtime = float(values.get("total", calc_runtime + plot_runtime + overhead_runtime))
            parts = f"{_fmt_tool(tool)} {calc_runtime:.2f}s calc / {plot_runtime:.2f}s plot"
            if overhead_runtime > 1e-3:
                parts += f" / {overhead_runtime:.2f}s overhead"
            runtime_chunks.append(f"{parts} ({total_runtime:.2f}s total)")
        runtime_para.text = "Runtime overview: " + ", ".join(runtime_chunks) + "."
        
        # Add timing note for benchmark methodology
        timing_note = frame.add_paragraph()
        timing_note.level = 1
        timing_note.font.size = Pt(12)
        timing_note.font.italic = True
        timing_note.text = (
            "Note: FastMDAnalysis shows single-line analyze() orchestrator with integrated file I/O. "
            "Calc/plot split estimated using 17.5%/82.5% ratio from individual benchmarks. "
            "Other tools show aggregate with file I/O separated. Memory includes all analyses running together."
        )

    if peak_mem_totals:
        mem_para = frame.add_paragraph()
        mem_para.level = 1
        mem_para.font.size = Pt(16)
        mem_chunks = []
        for tool, values in peak_mem_totals.items():
            calc_mem = float(values.get("calc", 0.0))
            plot_mem = float(values.get("plot", 0.0))
            total_mem = float(values.get("total", calc_mem + plot_mem))
            mem_chunks.append(
                f"{_fmt_tool(tool)} {calc_mem:.2f} MB calc / {plot_mem:.2f} MB plot ({total_mem:.2f} MB peak)"
            )
        mem_para.text = "Peak memory overview: " + ", ".join(mem_chunks) + "."

    closing = frame.add_paragraph()
    closing.level = 1
    closing.font.size = Pt(16)
    closing.text = "Charts on the following slides illustrate these touchpoint gaps."


def add_combined_overview_slide(prs: Presentation, overview_dir: Path) -> None:
    """Add a single slide containing four overview charts in a 2x2 grid.

    Charts (if present): external_modules.png, loc_totals.png,
    runtime_totals.png, peak_mem_totals.png
    """
    files = [
        ("External module dependencies", overview_dir / "external_modules.png"),
        ("Snippet lines of code (calc vs plot)", overview_dir / "loc_totals.png"),
        ("Runtime footprint overview", overview_dir / "runtime_totals.png"),
        ("Peak memory footprint overview", overview_dir / "peak_mem_totals.png"),
    ]
    # Only add the slide if at least one of the files exists
    if not any(p.exists() for _, p in files):
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9.0), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_run = title_frame.paragraphs[0].add_run()
    title_run.text = "Overview: Dependencies, LOC, Runtime, Memory"
    title_run.font.bold = True
    title_run.font.size = Pt(24)

    positions = [
        (Inches(0.5), Inches(1.0)),
        (Inches(5.0), Inches(1.0)),
        (Inches(0.5), Inches(4.2)),
        (Inches(5.0), Inches(4.2)),
    ]
    image_width = Inches(4.0)

    for (_caption, path), (left, top) in zip(files, positions):
        if not path.exists():
            continue
        slide.shapes.add_picture(str(path), left, top, width=image_width)


def add_aggregated_overview_slide(prs: Presentation, overview_dir: Path) -> None:
    """Add a single slide showing aggregated metrics (sum of individual benchmarks).
    
    This is different from the combined_overview_slide which uses orchestrator data.
    This slide shows the aggregated LOC as 8 lines (same as orchestrator but semantically
    different - aggregated means sum of individual runs, not single analyze() call).
    
    Charts: same 4 as combined overview, but title indicates aggregated metrics.
    """
    files = [
        ("External module dependencies", overview_dir / "external_modules.png"),
        ("Snippet lines of code (calc vs plot)", overview_dir / "loc_totals.png"),
        ("Runtime footprint overview", overview_dir / "runtime_totals.png"),
        ("Peak memory footprint overview", overview_dir / "peak_mem_totals.png"),
    ]
    # Only add the slide if at least one of the files exists
    if not any(p.exists() for _, p in files):
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9.0), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_run = title_frame.paragraphs[0].add_run()
    title_run.text = "Aggregated Overview: Dependencies, LOC, Runtime, Memory"
    title_run.font.bold = True
    title_run.font.size = Pt(24)

    # Add a note explaining aggregated metrics
    note_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9.0), Inches(0.25))
    note_frame = note_box.text_frame
    note_frame.clear()
    note_para = note_frame.paragraphs[0]
    note_para.text = "(Aggregated metrics: sum of 4 individual benchmark runs - RMSD, RMSF, RG, Cluster)"
    note_para.font.size = Pt(11)
    note_para.font.italic = True

    positions = [
        (Inches(0.5), Inches(1.05)),
        (Inches(5.0), Inches(1.05)),
        (Inches(0.5), Inches(4.25)),
        (Inches(5.0), Inches(4.25)),
    ]
    image_width = Inches(4.0)

    for (_caption, path), (left, top) in zip(files, positions):
        if not path.exists():
            continue
        slide.shapes.add_picture(str(path), left, top, width=image_width)


def _count_noncomment_lines(path: Path) -> int:
    if not path.exists():
        return 0
    count = 0
    for line in path.read_text(encoding="utf-8").splitlines():
        s = line.strip()
        if not s or s.startswith("#"):
            continue
        count += 1
    return count


def add_orchestrator_single_run_slide(prs: Presentation, results_root: Path, overview_dir: Path) -> None:
    """Create a slide that compares running all analyses individually (aggregated)
    vs running them in a single FastMDAnalysis .analyze(...) invocation.
    """
    orchestrator_dir = results_root / f"orchestrator_{DATASET_SLUG}"
    if not orchestrator_dir.exists():
        return

    instrumentation_json = overview_dir / "instrumentation_summary.json"
    try:
        summary = json.loads(instrumentation_json.read_text(encoding="utf-8")) if instrumentation_json.exists() else {}
    except Exception:
        summary = {}

    benchmarks = summary.get("metadata", {}).get("benchmarks", []) if summary else []
    if not benchmarks:
        benchmarks = ["rmsd", "rmsf", "rg", "cluster"]

    agg_calc_runtime = agg_plot_runtime = agg_over_runtime = 0.0
    agg_calc_mem = agg_plot_mem = agg_over_mem = 0.0
    agg_runtime = 0.0
    agg_peak = 0.0
    agg_loc_calc = agg_loc_plot = agg_loc_total = 0.0
    summary_runtime: dict | None = None
    summary_memory: dict | None = None
    summary_loc_calc = summary_loc_plot = summary_loc_total = None
    per_analysis_totals: dict[str, float] = {}
    if summary:
        runtime_totals = summary.get("runtime_totals", {})
        peak_mem_totals = summary.get("peak_mem_totals", {})
        loc_totals = summary.get("loc_totals", {})
        fd = "fastmdanalysis"
        summary_runtime = runtime_totals.get(fd)
        summary_memory = peak_mem_totals.get(fd)
        if loc_totals.get(fd):
            loc_vals = loc_totals[fd]
            summary_loc_calc = float(loc_vals.get("calc", 0))
            summary_loc_plot = float(loc_vals.get("plot", 0))
            summary_loc_total = float(
                loc_vals.get(
                    "total",
                    loc_vals.get("calc", 0) + loc_vals.get("plot", 0),
                )
            )
    for bench in benchmarks:
        bench_dir = results_root / f"{bench}_{DATASET_SLUG}"
        metrics_path = bench_dir / "metrics.json"
        if not metrics_path.exists():
            continue
        try:
            bench_metrics = json.loads(metrics_path.read_text(encoding="utf-8"))
        except Exception:
            continue
        summary_entry = bench_metrics.get("summary", {}).get("fastmdanalysis", {})
        calc_rt = float(summary_entry.get("calc_s", {}).get("mean", 0.0))
        plot_rt = float(summary_entry.get("plot_s", {}).get("mean", 0.0))
        total_rt = float(
            summary_entry.get("total_s", {}).get(
                "mean",
                summary_entry.get("elapsed_s", {}).get("mean", calc_rt + plot_rt),
            )
        )
        overhead_rt = max(total_rt - calc_rt - plot_rt, 0.0)
        agg_calc_runtime += calc_rt
        agg_plot_runtime += plot_rt
        agg_over_runtime += overhead_rt
        agg_runtime += total_rt
        per_analysis_totals[bench] = total_rt

        calc_mem = float(summary_entry.get("calc_mem_mb", {}).get("mean", 0.0))
        plot_mem = float(summary_entry.get("plot_mem_mb", {}).get("mean", 0.0))
        peak_mem = float(summary_entry.get("peak_mem_mb", {}).get("mean", max(calc_mem, plot_mem)))
        overhead_mem = max(peak_mem - calc_mem - plot_mem, 0.0)
        if peak_mem >= agg_peak:
            agg_peak = peak_mem
            agg_calc_mem = calc_mem
            agg_plot_mem = plot_mem
            agg_over_mem = overhead_mem

        loc_stats = summary_entry.get("loc", {})
        calc_loc = float(loc_stats.get("calc", 0.0) or 0.0)
        plot_loc = float(loc_stats.get("plot", 0.0) or 0.0)
        total_loc = float(loc_stats.get("total", calc_loc + plot_loc))
        if calc_loc or plot_loc or total_loc:
            agg_loc_calc += calc_loc
            agg_loc_plot += plot_loc
            agg_loc_total += total_loc

    runtime_calc_sum = agg_calc_runtime
    runtime_plot_sum = agg_plot_runtime
    runtime_over_sum = agg_over_runtime
    runtime_total_sum = agg_runtime

    memory_calc_sum = agg_calc_mem
    memory_plot_sum = agg_plot_mem
    memory_over_sum = agg_over_mem
    memory_total_sum = agg_peak

    if runtime_total_sum > 0.0:
        agg_calc_runtime = runtime_calc_sum
        agg_plot_runtime = runtime_plot_sum
        agg_over_runtime = max(runtime_total_sum - (agg_calc_runtime + agg_plot_runtime), runtime_over_sum, 0.0)
        agg_runtime = runtime_total_sum
    elif summary_runtime:
        agg_calc_runtime = float(summary_runtime.get("calc", 0.0))
        agg_plot_runtime = float(summary_runtime.get("plot", 0.0))
        agg_over_runtime = float(summary_runtime.get("overhead", 0.0))
        agg_runtime = float(summary_runtime.get("total", agg_calc_runtime + agg_plot_runtime + agg_over_runtime))
    else:
        agg_runtime = agg_calc_runtime + agg_plot_runtime + agg_over_runtime
        agg_over_runtime = max(agg_over_runtime, 0.0)

    if memory_total_sum > 0.0:
        agg_calc_mem = memory_calc_sum
        agg_plot_mem = memory_plot_sum
        agg_over_mem = max(memory_total_sum - (agg_calc_mem + agg_plot_mem), memory_over_sum, 0.0)
        agg_peak = memory_total_sum
    elif summary_memory:
        agg_calc_mem = float(summary_memory.get("calc", 0.0))
        agg_plot_mem = float(summary_memory.get("plot", 0.0))
        agg_over_mem = max(
            float(summary_memory.get("total", 0.0)) - (agg_calc_mem + agg_plot_mem),
            0.0,
        )
        agg_peak = float(summary_memory.get("total", agg_calc_mem + agg_plot_mem + agg_over_mem))
    else:
        agg_peak = agg_calc_mem + agg_plot_mem + agg_over_mem
        agg_over_mem = max(agg_over_mem, 0.0)

    if agg_loc_total == 0.0 and summary_loc_total is not None:
        agg_loc_calc = summary_loc_calc or 0.0
        agg_loc_plot = summary_loc_plot or 0.0
        agg_loc_total = summary_loc_total or (agg_loc_calc + agg_loc_plot)

    # Single-run orchestrator metrics
    orch_metrics_path = orchestrator_dir / "metrics.json"
    single_total_runtime = single_calc_runtime = single_plot_runtime = None
    single_peak = single_calc_mem = single_plot_mem = None
    single_loc_calc = single_loc_plot = single_loc_total = None
    single_analysis_breakdown: dict[str, float] = {}
    if orch_metrics_path.exists():
        try:
            orch = json.loads(orch_metrics_path.read_text(encoding="utf-8"))
            s = orch.get("summary", {}).get("fastmdanalysis", {})

            def _metric_mean(summary: dict, key: str) -> float | None:
                value = summary.get(key)
                if not isinstance(value, dict):
                    return None
                mean = value.get("mean")
                if mean is None:
                    return None
                try:
                    return float(mean)
                except (TypeError, ValueError):
                    return None

            single_total_runtime = _metric_mean(s, "total_s")
            if single_total_runtime is None:
                single_total_runtime = _metric_mean(s, "elapsed_s")
            single_calc_runtime = _metric_mean(s, "calc_s")
            single_plot_runtime = _metric_mean(s, "plot_s")
            single_peak = _metric_mean(s, "peak_mem_mb")
            single_calc_mem = _metric_mean(s, "calc_mem_mb")
            single_plot_mem = _metric_mean(s, "plot_mem_mb")
            per_analysis_stats = s.get("per_analysis_s", {})
            if isinstance(per_analysis_stats, dict):
                for name, stats in per_analysis_stats.items():
                    if isinstance(stats, dict):
                        try:
                            single_analysis_breakdown[name] = float(stats.get("mean", 0.0) or 0.0)
                        except (TypeError, ValueError):
                            continue
            loc_data = s.get("loc", {})
            if isinstance(loc_data, dict):
                single_loc_calc = float(loc_data.get("calc", loc_data.get("total", 0)))
                single_loc_plot = float(loc_data.get("plot", 0))
                single_loc_total = float(loc_data.get("total", single_loc_calc + single_loc_plot))
        except Exception:
            single_total_runtime = single_peak = None

    if single_loc_total is None:
        example_path = ROOT / "examples" / "api_analyze_all.py"
        example_lines = _count_noncomment_lines(example_path)
        if example_lines:
            single_loc_total = float(example_lines)

    if agg_runtime is None and single_total_runtime is None:
        return

    agg_calc_runtime = float(agg_calc_runtime or 0.0)
    agg_plot_runtime = float(agg_plot_runtime or 0.0)
    agg_over_runtime = float(agg_over_runtime or 0.0)
    agg_runtime = float(agg_runtime or (agg_calc_runtime + agg_plot_runtime + agg_over_runtime))

    runtime_metrics_missing = (
        single_calc_runtime is None
        and single_plot_runtime is None
        and single_total_runtime is None
    )
    if runtime_metrics_missing:
        single_calc_runtime_val = agg_calc_runtime
        single_plot_runtime_val = agg_plot_runtime
        single_runtime_total = agg_runtime
    else:
        single_calc_runtime_val = float(single_calc_runtime) if single_calc_runtime is not None else 0.0
        single_plot_runtime_val = float(single_plot_runtime) if single_plot_runtime is not None else 0.0
        if single_total_runtime is not None:
            single_runtime_total = float(single_total_runtime)
        else:
            single_runtime_total = single_calc_runtime_val + single_plot_runtime_val
    single_over_runtime = max(single_runtime_total - single_calc_runtime_val - single_plot_runtime_val, 0.0)

    agg_calc_mem = float(agg_calc_mem or 0.0)
    agg_plot_mem = float(agg_plot_mem or 0.0)
    agg_over_mem = float(agg_over_mem or 0.0)
    agg_peak = float(agg_peak or (agg_calc_mem + agg_plot_mem + agg_over_mem))

    memory_metrics_missing = (
        single_calc_mem is None
        and single_plot_mem is None
        and single_peak is None
    )
    if memory_metrics_missing:
        single_calc_mem_val = agg_calc_mem
        single_plot_mem_val = agg_plot_mem
        single_peak_val = agg_peak
    else:
        single_calc_mem_val = float(single_calc_mem) if single_calc_mem is not None else 0.0
        single_plot_mem_val = float(single_plot_mem) if single_plot_mem is not None else 0.0
        single_peak_val = float(single_peak) if single_peak is not None else (single_calc_mem_val + single_plot_mem_val)
    single_over_mem = max(single_peak_val - single_calc_mem_val - single_plot_mem_val, 0.0)

    # Aggregated view sums per-analysis runs; single-run metrics come directly from the orchestrated benchmark

    agg_loc_calc = float(agg_loc_calc or 0.0)
    agg_loc_plot = float(agg_loc_plot or 0.0)
    agg_loc_over = max(float(agg_loc_total or (agg_loc_calc + agg_loc_plot)) - agg_loc_calc - agg_loc_plot, 0.0)
    agg_loc_total = float(agg_loc_total or (agg_loc_calc + agg_loc_plot + agg_loc_over))

    single_loc_calc = float(single_loc_calc or 0.0)
    single_loc_plot = float(single_loc_plot or 0.0)
    single_loc_total = float(single_loc_total or (single_loc_calc + single_loc_plot))
    single_loc_over = max(single_loc_total - single_loc_calc - single_loc_plot, 0.0)

    labels = ["Aggregated (4 runs)", "Single analyze() run"]
    segment_order = ["Computation", "Plotting", "Overhead"]
    calc_color = "#1f77b4"
    plot_color = "#6baed6"
    overhead_color = "#7f7f7f"

    chart_specs = [
        (
            "Runtime breakdown",
            "Seconds",
            "s",
            {
                "Computation": [agg_calc_runtime, single_calc_runtime_val],
                "Plotting": [agg_plot_runtime, single_plot_runtime_val],
                "Overhead": [agg_over_runtime, single_over_runtime],
            },
            overview_dir / "orchestrator_single_run_runtime.png",
            {"loc": "upper left", "bbox_to_anchor": (0.0, 1.0)},
        ),
        (
            "Peak memory breakdown",
            "Megabytes",
            "MB",
            {
                "Computation": [agg_calc_mem, single_calc_mem_val],
                "Plotting": [agg_plot_mem, single_plot_mem_val],
                "Overhead": [agg_over_mem, single_over_mem],
            },
            overview_dir / "orchestrator_single_run_memory.png",
            {"loc": "upper right", "bbox_to_anchor": (1.0, 1.0)},
        ),
        (
            "Lines-of-code footprint",
            "Lines",
            "LOC",
            {
                "Computation": [agg_loc_calc, single_loc_calc],
                "Plotting": [agg_loc_plot, single_loc_plot],
                "Overhead": [agg_loc_over, single_loc_over],
            },
            overview_dir / "orchestrator_single_run_loc.png",
            {"loc": "upper right", "bbox_to_anchor": (1.0, 1.0)},
        ),
    ]

    def _save_stacked_chart(
        title: str,
        ylabel: str,
        unit: str,
        series_data: dict[str, list[float]],
        path: Path,
        legend_kwargs: dict[str, object] | None = None,
    ) -> None:
        fig, ax = plt.subplots(figsize=(7.0, 4.4))
        x = np.arange(len(labels))
        width = 0.45
        totals = np.zeros(len(labels))
        bottoms = np.zeros(len(labels))

        for segment in segment_order:
            values = np.array(series_data.get(segment, [0.0] * len(labels)), dtype=float)
            bars = ax.bar(
                x,
                values,
                width=width,
                bottom=bottoms,
                color={
                    "Computation": calc_color,
                    "Plotting": plot_color,
                    "Overhead": overhead_color,
                }[segment],
                label=segment,
            )
            for bar in bars:
                bar.set_edgecolor("#3a3a3a")
                bar.set_linewidth(0.8)
                if segment == "Plotting":
                    bar.set_hatch("//")
                    bar.set_alpha(0.7)
                elif segment == "Overhead":
                    bar.set_alpha(0.8)
                else:
                    bar.set_alpha(0.9)
            bottoms += values
            totals += values

        max_total = float(totals.max()) if len(totals) else 0.0

        ax.set_xticks(x)
        ax.set_xticklabels(labels, rotation=12)
        ax.set_ylabel(ylabel)
        ax.set_title(title)
        ax.grid(axis="y", alpha=0.3)
        ax.set_ylim(0, max_total * 1.25 + (0.1 if max_total < 1 else 0.0))

        default_legend = {"loc": "upper left", "bbox_to_anchor": (0.0, 1.0)}
        legend_opts = default_legend if legend_kwargs is None else {**default_legend, **legend_kwargs}
        ax.legend(**legend_opts)

        fig.tight_layout()
        fig.savefig(path, bbox_inches="tight")
        plt.close(fig)

    def _save_analysis_breakdown_chart(labels: list[str], aggregated_vals: list[float], single_vals: list[float], path: Path) -> None:
        if not labels:
            return
        fig, ax = plt.subplots(figsize=(7.0, 4.4))
        x = np.arange(len(labels))
        width = 0.38
        agg_bars = ax.bar(x - width / 2, aggregated_vals, width=width, color=calc_color, alpha=0.85, label="Aggregated runs")
        single_bars = ax.bar(x + width / 2, single_vals, width=width, color="#2ca02c", alpha=0.85, label="Single analyze()")

        ax.set_xticks(x)
        ax.set_xticklabels(labels, rotation=12)
        ax.set_ylabel("Seconds")
        ax.set_title("Per-analysis time distribution")
        ax.grid(axis="y", alpha=0.3)
        max_total = max((aggregated_vals or [0]) + (single_vals or [0]))
        ax.set_ylim(0, max_total * 1.25 + (0.1 if max_total < 1 else 0.0))
        ax.legend(loc="upper left", bbox_to_anchor=(0.0, 1.0))

        def _annotate(bars, data):
            for bar, value in zip(bars, data):
                if value <= 0:
                    continue
                ax.text(
                    bar.get_x() + bar.get_width() / 2,
                    bar.get_height() + max(0.05, value * 0.03),
                    f"{value:.2f}s",
                    ha="center",
                    va="bottom",
                    fontsize=10,
                )

        _annotate(agg_bars, aggregated_vals)
        _annotate(single_bars, single_vals)

        fig.tight_layout()
        fig.savefig(path, bbox_inches="tight")
        plt.close(fig)

    chart_paths: list[Path] = []
    for title, ylabel, unit, series_data, path, legend_opts in chart_specs:
        _save_stacked_chart(title, ylabel, unit, series_data, path, legend_kwargs=legend_opts)
        chart_paths.append(path)

    analysis_breakdown_path: Path | None = None
    if per_analysis_totals and single_analysis_breakdown:
        labels_map = {
            name: name.upper() if name.upper() != name else name
            for name in set(per_analysis_totals) | set(single_analysis_breakdown)
        }
        analysis_labels: list[str] = []
        aggregated_vals: list[float] = []
        single_vals: list[float] = []
        considered: set[str] = set()
        for bench in benchmarks:
            if bench in considered:
                continue
            agg_val = float(per_analysis_totals.get(bench, 0.0))
            single_val = float(single_analysis_breakdown.get(bench, 0.0))
            if agg_val == 0.0 and single_val == 0.0:
                considered.add(bench)
                continue
            analysis_labels.append(labels_map.get(bench, bench.upper()))
            aggregated_vals.append(agg_val)
            single_vals.append(single_val)
            considered.add(bench)
        for name, value in single_analysis_breakdown.items():
            if name in considered:
                continue
            agg_val = float(per_analysis_totals.get(name, 0.0))
            single_val = float(value or 0.0)
            if agg_val == 0.0 and single_val == 0.0:
                continue
            analysis_labels.append(labels_map.get(name, name.upper()))
            aggregated_vals.append(agg_val)
            single_vals.append(single_val)
        if analysis_labels and (max(aggregated_vals + single_vals) > 0.0):
            analysis_breakdown_path = overview_dir / "orchestrator_single_run_analysis_breakdown.png"
            _save_analysis_breakdown_chart(analysis_labels, aggregated_vals, single_vals, analysis_breakdown_path)

    runtime_chart_path, memory_chart_path, loc_chart_path = chart_paths

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if analysis_breakdown_path and analysis_breakdown_path.exists():
        positions = [
            (runtime_chart_path, Inches(0.5), Inches(1.1), Inches(4.3)),
            (memory_chart_path, Inches(5.0), Inches(1.1), Inches(4.3)),
            (analysis_breakdown_path, Inches(0.5), Inches(4.9), Inches(4.3)),
            (loc_chart_path, Inches(5.0), Inches(4.9), Inches(4.3)),
        ]
    else:
        positions = [
            (runtime_chart_path, Inches(0.5), Inches(1.1), Inches(4.6)),
            (memory_chart_path, Inches(5.0), Inches(1.1), Inches(4.6)),
            (loc_chart_path, Inches(0.5), Inches(4.9), Inches(9.1)),
        ]
    for chart_path, left, top, width in positions:
        if chart_path.exists():
            slide.shapes.add_picture(str(chart_path), left, top, width=width)

    caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.6))
    caption_frame = caption_box.text_frame
    caption_frame.clear()
    caption_run = caption_frame.paragraphs[0].add_run()
    caption_run.text = "Single-run vs aggregated FastMDAnalysis"
    caption_run.font.size = Pt(24)
    caption_run.font.bold = True

    note_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9.0), Inches(1.0))
    note_frame = note_box.text_frame
    note_frame.clear()
    note_para = note_frame.paragraphs[0]
    note_text = (
        "Aggregated bars sum the computation, plotting, and overhead from individual analyze runs."
        " Single-run bars reflect the orchestrated analyze(...) benchmark metrics for those same slices."
    )
    if analysis_breakdown_path and analysis_breakdown_path.exists():
        note_text += " Bottom-left chart breaks down runtime per analysis for the single run (slides counted as plotting)."
    note_para.text = note_text
    note_para.font.size = Pt(12)
    note_para.font.italic = True


def add_instrumentation_combined_slide(prs: Presentation, overview_dir: Path) -> None:
    """Add a single slide containing instrumentation overview, modules/functions per benchmark
    and the orchestrator LOC advantage chart in a 2x2 grid.

    Files (if present): instrumentation_overview.png, modules_per_benchmark.png,
    functions_per_benchmark.png, orchestrator_loc_advantage.png
    """
    files = [
        ("Instrumentation overview", overview_dir / "instrumentation_overview.png"),
        ("Modules touched per benchmark", overview_dir / "modules_per_benchmark.png"),
        ("Functions touched per benchmark", overview_dir / "functions_per_benchmark.png"),
        ("Orchestrator LOC advantage", overview_dir / "orchestrator_loc_advantage.png"),
    ]
    if not any(p.exists() for _, p in files):
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9.0), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_run = title_frame.paragraphs[0].add_run()
    title_run.text = "Instrumentation overview: touchpoints, modules, functions, LOC"
    title_run.font.bold = True
    title_run.font.size = Pt(24)

    positions = [
        (Inches(0.5), Inches(1.0)),
        (Inches(5.0), Inches(1.0)),
        (Inches(0.5), Inches(4.2)),
        (Inches(5.0), Inches(4.2)),
    ]
    image_width = Inches(4.0)

    for (_caption, path), (left, top) in zip(files, positions):
        if not path.exists():
            continue
        slide.shapes.add_picture(str(path), left, top, width=image_width)


def build_presentation(output_path: Path) -> None:
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_text = "FastMdAnalysis Benchmarks"
    if DATASET_LABEL:
        title_text += f" – {DATASET_LABEL}"
    title_slide.shapes.title.text = title_text
    subtitle = title_slide.placeholders[1]
    subtitle_text = (
        "Runtime, memory, and usability comparisons versus MDTraj and MDAnalysis\n"
        "Note: third-party metrics include helper scripts for exporting data and plots."
    )
    if DATASET_LABEL:
        subtitle_text += f"\nDataset: {DATASET_LABEL}"
    subtitle.text = subtitle_text

    # Add a combined overview slide containing the four aggregate charts
    # (external modules, LOC totals, runtime totals, peak mem totals).
    try:
        add_combined_overview_slide(prs, OVERVIEW_DIR)
    except Exception:
        # Non-fatal: if something goes wrong adding the combined slide,
        # continue building the rest of the presentation.
        pass

    # Add aggregated overview slide (sum of individual benchmarks, not orchestrator)
    try:
        add_aggregated_overview_slide(prs, OVERVIEW_DIR)
    except Exception:
        # Non-fatal: continue if something fails while adding aggregated slide
        pass

    # Add a combined instrumentation slide containing instrumentation_overview,
    # modules_per_benchmark, functions_per_benchmark, and orchestrator_loc_advantage
    try:
        add_instrumentation_combined_slide(prs, OVERVIEW_DIR)
    except Exception:
        # Non-fatal: continue if something fails while adding instrumentation slide
        pass

    benchmark_dirs = list(iter_benchmark_dirs(RESULTS_ROOT, DATASET_SLUG))
    if not benchmark_dirs:
        raise FileNotFoundError(
            f"No benchmark outputs found for dataset '{DATASET_SLUG}'. Run the benchmark scripts first."
        )
    summary: dict | None = None
    if INSTRUMENTATION_JSON.exists():
        summary = json.loads(INSTRUMENTATION_JSON.read_text(encoding="utf-8"))

    tools: List[str] = []
    aggregated: dict = {}
    aggregated_details: dict = {}
    per_benchmark = {}

    def _fmt_tool(name: str) -> str:
        # Normalize display names for the primary tools we benchmark.
        mapping = {
            "fastmdanalysis": "FastMdAnalysis",
            "mdtraj": "MdTraj",
            "mdanalysis": "MdAnalysis",
        }
        return mapping.get(name.lower(), name.capitalize())

    def _fmt_range(values: List[int]) -> str:
        if not values:
            return "0"
        low = min(values)
        high = max(values)
        return str(low) if low == high else f"{low}-{high}"

    if summary:
        tools = summary.get("metadata", {}).get("tools", [])
        aggregated = summary.get("aggregated", {})
        aggregated_details = summary.get("aggregated_details", {})
        per_benchmark = summary.get("per_benchmark", {})
        summary_dataset = summary.get("metadata", {}).get("dataset")
        if summary_dataset and summary_dataset != DATASET_SLUG:
            raise ValueError(
                f"Instrumentation summary dataset '{summary_dataset}' does not match selected dataset '{DATASET_SLUG}'."
            )

    # Add per-benchmark slides (charts etc.) but skip generating individual
    # per-benchmark instrumentation overview slides — we'll include the
    # instrumentation overview, modules/functions per benchmark, and the
    # orchestrator LOC advantage on a single combined instrumentation slide.
    for benchmark_dir in benchmark_dirs:
        add_benchmark_slide(prs, benchmark_dir)

    # Create an orchestrator LOC advantage chart comparing the LOC required
    # to run all analyses per-tool (assumes non-orchestrated tools need one
    # snippet per analysis) versus FastMDAnalysis which can orchestrate all
    # analyses with a single snippet.
    if summary:
        try:
            # Use aggregated_loc_totals which has the summed LOC from individual benchmarks
            aggregated_loc = summary.get("aggregated_loc_totals", {})
            orchestrator_loc = summary.get("orchestrator_loc_totals", {})
            loc_totals = summary.get("loc_totals", {})
            
            ordered_tools = tools or list(loc_totals.keys())
            analyses_count = 4
            # Display/name mapping and standard colors requested
            name_map = {
                "fastmdanalysis": "FastMdAnalysis",
                "mdtraj": "MdTraj",
                "mdanalysis": "MdAnalysis",
            }
            color_map = {
                "fastmdanalysis": "#1f77b4",  # blue
                "mdtraj": "#ff7f0e",         # orange
                "mdanalysis": "#7f7f7f",     # gray
            }

            labels = []
            aggregate_loc = []
            orchestrator_vals = []
            
            for tool in ordered_tools:
                # Use aggregated_loc_totals which contains the summed LOC from individual benchmarks
                vals = aggregated_loc.get(tool, loc_totals.get(tool, {}))
                total_loc = int(vals.get("calc", 0) + vals.get("plot", 0))
                
                labels.append(name_map.get(tool.lower(), tool.capitalize()))
                aggregate_loc.append(total_loc)
                
                # For FastMDAnalysis, add orchestrator LOC
                if tool.lower() == "fastmdanalysis" and orchestrator_loc.get(tool):
                    orch_vals = orchestrator_loc.get(tool, {})
                    orch_total = int(orch_vals.get("calc", 0) + orch_vals.get("plot", 0))
                    orchestrator_vals.append(orch_total)
                else:
                    orchestrator_vals.append(total_loc)

            # Plot and save - show comparison between aggregate and orchestrator
            fig, ax = plt.subplots(figsize=(8, 4.2))
            x = np.arange(len(labels))
            width = 0.35
            colors_list = [color_map.get(tool.lower(), "#bdbdbd") for tool in ordered_tools]
            
            # For tools without orchestrator, show single bar
            # For FastMDAnalysis, show two bars
            bars_aggregate = ax.bar(x - width/2, aggregate_loc, width, label='Aggregate (4 separate calls)', 
                                   color=colors_list, alpha=0.7)
            bars_orchestrator = ax.bar(x + width/2, orchestrator_vals, width, label='Orchestrator (analyze)', 
                                      color=colors_list)
            
            ax.set_xticks(x)
            ax.set_xticklabels(labels)
            ax.set_ylabel("Lines of Code")
            ax.set_title("LOC Comparison: Aggregate vs Orchestrator (4 analyses)")
            ax.legend()
            
            # Annotate bars with values
            for bars in [bars_aggregate, bars_orchestrator]:
                for bar in bars:
                    height = bar.get_height()
                    if height > 0:
                        ax.text(bar.get_x() + bar.get_width()/2., height,
                               f'{int(height)}',
                               ha='center', va='bottom', fontsize=9)
            
            fig.tight_layout()
            outpath = OVERVIEW_DIR / "orchestrator_loc_advantage.png"
            fig.savefig(outpath, bbox_inches="tight")
            plt.close(fig)

            # Orchestrator chart: do not add an individual slide here. It will
            # be included on the combined instrumentation slide instead.
        except Exception as e:
            # Non-fatal; continue building presentation
            print(f"Warning: Failed to create orchestrator LOC advantage chart: {e}")
            pass

    if summary:
        # Add a slide comparing single-run orchestrator vs aggregated separate runs
        try:
            add_orchestrator_single_run_slide(prs, RESULTS_ROOT, OVERVIEW_DIR)
        except Exception:
            pass

        add_summary_slide(prs, summary)

        instrumentation_notes: List[str] | None = None
        if INSTRUMENTATION_CHART.exists() and aggregated:
            instrumentation_notes = []
            ordered_tools = tools or list(aggregated.keys())
            for tool in ordered_tools:
                stats = aggregated.get(tool, {})
                instrumentation_notes.append(
                    f"{_fmt_tool(tool)}: {stats.get('modules', 0)} modules, {stats.get('functions', 0)} functions, "
                    f"{stats.get('external_modules', 0)} external imports, "
                    f"{stats.get('attempts', 0)} runs, {stats.get('successes', 0)} successful, {stats.get('exceptions', 0)} exceptions"
                )
            loc_totals = summary.get("loc_totals", {})
            if loc_totals:
                footprint_chunks = []
                for tool, vals in loc_totals.items():
                    calc_lines = vals.get("calc", 0)
                    plot_lines = vals.get("plot", 0)
                    total_lines = calc_lines + plot_lines
                    footprint_chunks.append(
                        f"{_fmt_tool(tool)} {calc_lines} calc / {plot_lines} plot ({total_lines} total) lines"
                    )
                footprint = ", ".join(footprint_chunks)
                instrumentation_notes.append(f"Snippet footprint: {footprint}")
            runtime_totals = summary.get("runtime_totals", {})
            if runtime_totals:
                runtime_chunks = []
                for tool, values in runtime_totals.items():
                    calc_runtime = float(values.get("calc", 0.0))
                    plot_runtime = float(values.get("plot", 0.0))
                    overhead_runtime = float(values.get("overhead", 0.0))
                    total_runtime = float(values.get("total", calc_runtime + plot_runtime + overhead_runtime))
                    parts = f"{_fmt_tool(tool)} {calc_runtime:.2f}s calc / {plot_runtime:.2f}s plot"
                    if overhead_runtime > 1e-3:
                        parts += f" / {overhead_runtime:.2f}s overhead"
                    runtime_chunks.append(f"{parts} ({total_runtime:.2f}s total)")
                instrumentation_notes.append("Aggregate runtime: " + ", ".join(runtime_chunks))
            peak_mem_totals = summary.get("peak_mem_totals", {})
            if peak_mem_totals:
                mem_chunks = []
                for tool, values in peak_mem_totals.items():
                    calc_mem = float(values.get("calc", 0.0))
                    plot_mem = float(values.get("plot", 0.0))
                    total_mem = float(values.get("total", calc_mem + plot_mem))
                    mem_chunks.append(
                        f"{_fmt_tool(tool)} {calc_mem:.2f} MB calc / {plot_mem:.2f} MB plot ({total_mem:.2f} MB peak)"
                    )
                instrumentation_notes.append("Aggregate peak memory: " + ", ".join(mem_chunks))
        # Do not add a standalone instrumentation overview slide here; the
        # instrumentation overview image is included on the combined
        # instrumentation slide along with modules/functions and the
        # orchestrator LOC advantage.

    for chart in _overview_charts():
        chart_path = chart["path"]
        if not chart_path.exists():
            continue

        notes: List[str] | None = None
        detail_field = chart.get("detail_field")
        if detail_field and tools and aggregated and aggregated_details:
            notes = []
            for tool in tools:
                stats = aggregated.get(tool, {})
                detail_items = aggregated_details.get(tool, {}).get(detail_field, [])
                amount = stats.get(detail_field, len(detail_items))
                joined = ", ".join(detail_items) if detail_items else "(none)"
                notes.append(f"{_fmt_tool(tool)} ({amount}): {joined}")

        summary_key = chart.get("summary_key")
        if summary_key == "loc_totals":
            loc_totals = summary.get("loc_totals", {}) if summary else {}
            if loc_totals:
                notes = notes or []
                fast = loc_totals.get("fastmdanalysis", {})
                if fast:
                    fast_calc = fast.get("calc", 0)
                    fast_plot = fast.get("plot", 0)
                    other_plot = [vals.get("plot", 0) for name, vals in loc_totals.items() if name != "fastmdanalysis"]
                    other_calc = [vals.get("calc", 0) for name, vals in loc_totals.items() if name != "fastmdanalysis"]
                    if other_plot:
                        notes.append(
                            "FastMDAnalysis keeps plotting helpers lean "
                            f"({fast_plot} lines vs {_fmt_range(other_plot)} elsewhere)."
                        )
                    if other_calc:
                        notes.append(
                            "Calculation snippets stay compact too "
                            f"({fast_calc} lines vs {_fmt_range(other_calc)} for the others)."
                        )
                for tool in tools or []:
                    values = loc_totals.get(tool, {})
                    if values:
                        calc_lines = values.get("calc", 0)
                        plot_lines = values.get("plot", 0)
                        total_lines = calc_lines + plot_lines
                        notes.append(
                            f"{_fmt_tool(tool)}: {calc_lines} calc lines, {plot_lines} plotting lines, {total_lines} total lines"
                        )
        elif summary_key == "runtime_totals" and summary:
            runtime_totals = summary.get("runtime_totals", {})
            if runtime_totals:
                notes = notes or []
                for tool in tools or []:
                    values = runtime_totals.get(tool, {})
                    if not values:
                        continue
                    calc_runtime = float(values.get("calc", 0.0))
                    plot_runtime = float(values.get("plot", 0.0))
                    overhead_runtime = float(values.get("overhead", 0.0))
                    total_runtime = float(values.get("total", calc_runtime + plot_runtime + overhead_runtime))
                    parts = f"{_fmt_tool(tool)}: {calc_runtime:.2f}s calc, {plot_runtime:.2f}s plot"
                    if overhead_runtime > 1e-3:
                        parts += f", {overhead_runtime:.2f}s overhead"
                    parts += f" ({total_runtime:.2f}s total)"
                    notes.append(parts)
        elif summary_key == "peak_mem_totals" and summary:
            peak_mem_totals = summary.get("peak_mem_totals", {})
            if peak_mem_totals:
                notes = notes or []
                for tool in tools or []:
                    values = peak_mem_totals.get(tool, {})
                    if not values:
                        continue
                    calc_mem = float(values.get("calc", 0.0))
                    plot_mem = float(values.get("plot", 0.0))
                    total_mem = float(values.get("total", calc_mem + plot_mem))
                    notes.append(
                        f"{_fmt_tool(tool)}: {calc_mem:.2f} MB calc, {plot_mem:.2f} MB plot ({total_mem:.2f} MB peak)"
                    )

        # Skip individual slides for charts that are included in the combined
        # overview slides (we already added them earlier). This includes both
        # the aggregate overview set and the instrumentation combined set.
        combined_files = {
            "external_modules.png",
            "loc_totals.png",
            "runtime_totals.png",
            "peak_mem_totals.png",
            "instrumentation_overview.png",
            "modules_per_benchmark.png",
            "functions_per_benchmark.png",
            "orchestrator_loc_advantage.png",
        }
        if chart_path.name in combined_files:
            continue

        add_chart_slide(
            prs,
            chart["title"],
            chart_path,
            chart.get("caption", ""),
            notes=notes,
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    if output_path.exists():
        output_path.unlink()
    prs.save(output_path)
    print(f"Saved benchmark slideshow to {output_path}")


def main(argv: list[str] | None = None) -> None:
    parser = argparse.ArgumentParser(description="Generate FastMDAnalysis benchmark slideshow.")
    available_datasets = sorted({name.lower() for name in list_datasets()})
    parser.add_argument(
        "--dataset",
        default=DATASET_SLUG,
        type=str.lower,
        choices=available_datasets,
        help="Dataset identifier whose benchmark results should be included.",
    )
    parser.add_argument(
        "--output",
        type=Path,
        help="Optional output path for the generated PowerPoint.",
    )
    args = parser.parse_args(argv)

    _set_dataset(args.dataset, output_path=args.output)
    if not OVERVIEW_DIR.exists():
        raise FileNotFoundError(
            f"Overview directory '{OVERVIEW_DIR}' not found for dataset '{DATASET_SLUG}'. "
            "Run aggregate_instrumentation.py for this dataset first."
        )

    build_presentation(OUTPUT_PPTX)


if __name__ == "__main__":
    main()
