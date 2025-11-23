from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
BENCHMARK_DIR = ROOT / "assets" / "benchmarks"
PRESENTATION_PATH = ROOT / "benchmark_presentation.pptx"

BENCHMARK_COMMAND = (
    "python benchmark_scaling.py --frames 500 1000 2000 5000 --iterations 1\n"
    "python scripts/plot_scaling.py --mode both --combined-y-scale log linear"
)

BENCHMARK_SLIDES = [
    ("Benchmark Scaling — Runtime (Log)", BENCHMARK_DIR / "combined_runtime_scaling_log.png"),
    ("Benchmark Scaling — Runtime (Linear)", BENCHMARK_DIR / "combined_runtime_scaling_linear.png"),
    ("Benchmark Scaling — Memory (Log)", BENCHMARK_DIR / "combined_memory_scaling_log.png"),
    ("Benchmark Scaling — Memory (Linear)", BENCHMARK_DIR / "combined_memory_scaling_linear.png"),
    ("FastMDAnalysis – Runtime Scaling", BENCHMARK_DIR / "fastmdanalysis_runtime_scaling.png"),
    ("FastMDAnalysis – Memory Scaling", BENCHMARK_DIR / "fastmdanalysis_memory_scaling.png"),
    ("MDTraj – Runtime Scaling", BENCHMARK_DIR / "mdtraj_runtime_scaling.png"),
    ("MDTraj – Memory Scaling", BENCHMARK_DIR / "mdtraj_memory_scaling.png"),
    ("MDAnalysis – Runtime Scaling", BENCHMARK_DIR / "mdanalysis_runtime_scaling.png"),
    ("MDAnalysis – Memory Scaling", BENCHMARK_DIR / "mdanalysis_memory_scaling.png"),
]


def add_benchmark_slides() -> Path:
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for title, image_path in BENCHMARK_SLIDES:
        if not image_path.exists():
            print(f"[warn] Missing image for slide '{title}': {image_path}")
            continue
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image_path),
            left=Inches(0.6),
            top=Inches(0.6),
            width=Inches(8.6),
        )
        textbox = slide.shapes.add_textbox(Inches(0.6), Inches(6.4), Inches(8.6), Inches(1.4))
        frame = textbox.text_frame
        frame.word_wrap = True

        title_run = frame.paragraphs[0].add_run()
        title_run.text = title
        title_run.font.size = Pt(20)
        title_run.font.bold = True

        command_para = frame.add_paragraph()
        command_run = command_para.add_run()
        command_run.text = BENCHMARK_COMMAND
        command_run.font.name = "Consolas"
        command_run.font.size = Pt(14)

    prs.save(PRESENTATION_PATH)
    return PRESENTATION_PATH


if __name__ == "__main__":
    output = add_benchmark_slides()
    print(f"Updated benchmark slideshow: {output}")
