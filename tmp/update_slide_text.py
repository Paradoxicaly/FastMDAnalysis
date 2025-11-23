from pptx import Presentation

old_text = "Ubiquitin Q99 scaling benchmark (frames 0:-1:10, 3 iterations per point)"
new_text = "Ubiquitin Q99 scaling benchmark (pure computation, frames 0:-1:10, 3 iterations per point)"

prs = Presentation('benchmark_presentation.pptx')

def set_text(shape, text):
    tf = shape.text_frame
    tf.clear()
    tf.text = text

for slide in prs.slides:
    for shape in slide.shapes:
        if not hasattr(shape, 'text'):
            continue
        if shape.text.strip() == old_text:
            set_text(shape, new_text)

prs.save('benchmark_presentation.pptx')
